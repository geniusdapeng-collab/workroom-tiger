#!/usr/bin/env python3
"""build-capability-pptx.py · 能力导览 PPT 生成器（消费 docs/capabilities.auto.json）
用法：python3 scripts/build-capability-pptx.py  → 产出 docs/capability-tour.pptx
数据由 generate-capabilities.mjs 生成，先跑 `pnpm capabilities`。
"""
import json, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = os.path.join(os.path.dirname(__file__), "..")
DATA = json.load(open(os.path.join(ROOT, "docs/capabilities.auto.json"), encoding="utf-8"))

CANDY_BG = RGBColor(0xFF, 0xF5, 0xF7)      # 蜜桃雾底
CANDY_MAIN = RGBColor(0xD4, 0x00, 0x2A)    # 深珊瑚主色
CANDY_TEXT = RGBColor(0x4A, 0x2B, 0x33)    # 深莓正文
CANDY_SUB = RGBColor(0x8A, 0x4B, 0x5A)     # 雾莓次级

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = CANDY_BG
    return s

def text(s, x, y, w, h, lines):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (t, size, color, bold, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = t
        f = r.font; f.size = Pt(size); f.color.rgb = color; f.bold = bold
        f.name = "PingFang SC"
    return tb

# ① 封面
s = slide()
text(s, 1, 2.2, 11.3, 3, [
    (f"{DATA['repo']} · 能力导览", 44, CANDY_MAIN, True, PP_ALIGN.CENTER),
    (DATA["description"], 20, CANDY_TEXT, False, PP_ALIGN.CENTER),
    (f"自动生成 · {DATA['generatedAt'][:10]} · 与代码同步", 14, CANDY_SUB, False, PP_ALIGN.CENTER),
])

# ② 5 分钟体验路径
s = slide()
text(s, 0.8, 0.5, 11.7, 1, [("🚀 5 分钟体验路径", 32, CANDY_MAIN, True, PP_ALIGN.LEFT)])
text(s, 0.8, 1.6, 11.7, 1, [("pnpm install && pnpm preview:all —— 无需真实后端/密钥，Mock 数据已固化", 18, CANDY_TEXT, False, PP_ALIGN.LEFT)])
rows = [("🖥 PC · B 端工作台", "http://localhost:3000", "经营主页 / 晨报 / 待审批 / 一句话目标"),
        ("📱 B 端移动", "http://localhost:3001", "高保真演示页 + 手机壳容器"),
        ("📱 C 端 AI 服务前台", "http://localhost:3002", "免登对话 / 服务 / 工单 / 消息 / 我的")]
for i, (a, b, c) in enumerate(rows):
    text(s, 1.0, 2.8 + i * 1.3, 11.3, 1.1, [(a, 22, CANDY_MAIN, True, PP_ALIGN.LEFT),
                                            (f"{b}    {c}", 16, CANDY_SUB, False, PP_ALIGN.LEFT)])

# ③ 每组一页
for g in DATA["groups"]:
    s = slide()
    text(s, 0.8, 0.5, 11.7, 1, [(f"{g['icon']} {g['title']}", 30, CANDY_MAIN, True, PP_ALIGN.LEFT)])
    lines = []
    for it in g["items"][:10]:
        lines.append((f"• {it['name']}", 17, CANDY_TEXT, True, PP_ALIGN.LEFT))
        lines.append((f"   {it['desc']} — {it['how']}", 13, CANDY_SUB, False, PP_ALIGN.LEFT))
    text(s, 0.9, 1.6, 11.6, 5.6, lines)

# ④ 结尾
s = slide()
text(s, 1, 2.6, 11.3, 2.4, [
    ("下一步", 32, CANDY_MAIN, True, PP_ALIGN.CENTER),
    ("二次开发读 AGENTS.md → pnpm agent:tour → docs/capability-map.md", 18, CANDY_TEXT, False, PP_ALIGN.CENTER),
    ("发布前 pnpm release:gate 全过（硬性门禁）", 16, CANDY_SUB, False, PP_ALIGN.CENTER),
])

out = os.path.join(ROOT, "docs/capability-tour.pptx")
prs.save(out)
print(f"✓ {out}（{len(prs.slides.__iter__.__self__._sldIdLst)} 页）")
